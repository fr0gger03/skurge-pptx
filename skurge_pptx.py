#!/usr/bin/env python3

import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt

###########
# ingest payload
###########

vPayloadArray={
  "Customer": "Acme Industries",
  "SA": "SAFirstname SALastname",
  "solution": "VMware Cloud on AWS",
  "tvType": "VMW/AWS funded Pilot",
  "services": [
    "NA"
  ],
  "estDate": "2022-01-10",
  "useCases": [
    "Data Center Extension - Test/Dev",
    "Data Center Extension - Footprint Expansion",
    "Data Center Extension - On-demand Capacity"
  ],
  "Tests": [
    {
      "RptTestCaseID": "6001",
      "RptTestService": "VMware Cloud on AWS",
      "RptTestCase": "Provision cloud SDDC",
      "RptResults": "Successful instantiation / creation of VMware Cloud on AWS Software Defined Datacenter (SDDC);  if using VCDR, provisioned through VCDR interface\n",
      "RptRequirements": "Vmware Account activated\nAWS Account activated\nAWS VPC created with necessary CIDR and subnets\nSDDC created and connected to AWS VPC"
    },
    {
      "RptTestCaseID": "6002",
      "RptTestService": "VMware Cloud on AWS",
      "RptTestCase": "Basic connectivity",
      "RptResults": "Communications between VMC/On-Prem – Direct Connect or IPSec VPN",
      "RptRequirements": "L3 VPN successfully established – OR - Direct Connect successfully established\nOpen / interact with vCenter in VMC \nCold migration of VM is successful"
    },
    {
      "RptTestCaseID": "6004",
      "RptTestService": "VMware Cloud on AWS",
      "RptTestCase": "Application Load testing - performance benchmark",
      "RptResults": "Compare performance of CUSTOMER application on-prem to in VMC on AWS",
      "RptRequirements": "Move / migrate CUSTOMER application to VMC on AWS\nGenerate load Compare results"
    }
  ]
}


###########
# instantiate variables for data ingest
###########
customer = vPayloadArray["Customer"]
setSA = vPayloadArray["SA"]
tvType = vPayloadArray["tvType"]
estDate = vPayloadArray["estDate"]
solution = vPayloadArray["solution"]
services = vPayloadArray["services"]
use_cases = vPayloadArray["useCases"]
success_count = len(vPayloadArray["Tests"][0])

###########
# open sample presentation
###########
prs = Presentation('sample.pptx')

###########
# instantiate slide layouts
###########
slide_layout_title = prs.slide_layouts[0]
slide_layout_section = prs.slide_layouts[1]
slide_layout_content = prs.slide_layouts[2]
slide_layout_titlesub = prs.slide_layouts[3]
slide_layout_2content = prs.slide_layouts[4]
slide_layout_blank = prs.slide_layouts[5]
slide_layout_thanks = prs.slide_layouts[6]

# to add a slide use sample below
# slide = prs.slides.add_slide(slide_layout_title)

###########
# update title slide
###########
slide_title = prs.slides[0]

# for shape in slide_title.placeholders:
#     print('%d %s' % (shape.placeholder_format.idx, shape.name))

shape_title = slide_title.placeholders[0]
shape_title.text = customer
shape_tvtype = slide_title.placeholders[10]
shape_tvtype.text = tvType
shape_sa = slide_title.placeholders[11]
shape_sa.text = setSA

###########
# update just the facts slide
###########
slide_jtf = prs.slides[1]

shape_content = slide_jtf.placeholders[14]
table = shape_content.table
cell_cust = table.cell(0, 1)
cell_estdate = table.cell(1, 1)
cell_tvtype = table.cell(2, 1)
cell_sa = table.cell(3, 1)
cell_uses = table.cell(4, 1)
cell_solution = table.cell(5, 1)
cell_services = table.cell(6, 1)

cell_cust.text = customer
cell_estdate.text = setSA
cell_tvtype.text = tvType
cell_sa.text = estDate
cell_uses.text = "\n".join(use_cases)
cell_solution.text = solution
cell_services.text = "\n".join(services)

###########
# update success criteria slide
###########
# slide_success = prs.slides.add_slide(slide_layout_content)
slide_success = prs.slides.add_slide(slide_layout_titlesub)
shape_title = slide_success.placeholders[0]
shape_title.text = "Success Criteria Overview"
shape_subtitle = slide_success.placeholders[10]
shape_subtitle.text = "Overview"
x, y, cx, cy = Inches(0.67), Inches(1.25), Inches(12), Inches(5)
shape_content = slide_success.shapes.add_table(success_count+1, 5, x, y, cx, cy)
table = shape_content.table

# populate table
cell_hdr1 = table.cell(0, 0)
cell_hdr2 = table.cell(0, 1)
cell_hdr3 = table.cell(0, 2)
cell_hdr4 = table.cell(0, 3)
cell_hdr5 = table.cell(0, 4)

cell_hdr1.text = "Test ID"
cell_hdr2.text = "Solution"
cell_hdr3.text = "Test"
cell_hdr4.text = "Expected results"
cell_hdr5.text = "Requirement mapping"

# iterate through success criteria, populate table
# reset case_counter
row_counter = 1
for case in vPayloadArray["Tests"]:
    case_testid = case.get("RptTestCaseID")
    cell_testid =  table.cell(row_counter, 0)
    cell_testid.text = case_testid
    case_test_solution = case.get("RptTestService")
    cell_test =  table.cell(row_counter, 1)
    cell_test.text = case_test_solution
    case_test = case.get("RptTestCase")
    cell_test =  table.cell(row_counter, 2)
    cell_test.text = case_test
    case_results = case.get("RptResults")
    cell_results =  table.cell(row_counter, 3)
    cell_results.text = case_results
    case_reqs = case.get("RptRequirements")
    cell_reqs =  table.cell(row_counter, 4)
    cell_reqs.text = case_reqs
    row_counter += 1

# force font size to 10
def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

for cell in iter_cells(table):
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(10)

# add one slide per test
for case in vPayloadArray["Tests"]:
    slide_test = prs.slides.add_slide(slide_layout_titlesub)
    x, y, cx, cy = Inches(0.67), Inches(1.25), Inches(12), Inches(5)
    shape_content = slide_test.shapes.add_table(3, 5, x, y, cx, cy)
    table = shape_content.table

    cell_hdr1 = table.cell(0, 0)
    cell_hdr2 = table.cell(0, 1)
    cell_hdr3 = table.cell(0, 2)
    cell_hdr4 = table.cell(0, 3)
    cell_hdr5 = table.cell(0, 4)

    cell_hdr1.text = "Test ID"
    cell_hdr2.text = "Solution"
    cell_hdr3.text = "Test"
    cell_hdr4.text = "Expected results"
    cell_hdr5.text = "Requirement mapping"

    case_testid = case.get("RptTestCaseID")
    cell_testid =  table.cell(1, 0)
    cell_testid.text = case_testid
    case_test_solution = case.get("RptTestService")
    cell_test =  table.cell(1, 1)
    cell_test.text = case_test_solution
    case_test = case.get("RptTestCase")
    cell_test =  table.cell(1, 2)
    cell_test.text = case_test
    case_results = case.get("RptResults")
    cell_results =  table.cell(1, 3)
    cell_results.text = case_results
    case_reqs = case.get("RptRequirements")
    cell_reqs =  table.cell(1, 4)
    cell_reqs.text = case_reqs
    shape_title = slide_test.placeholders[0]
    shape_title.text = "Success Criteria Detail"
    shape_subtitle = slide_test.placeholders[10]
    shape_subtitle.text = case_test
    # force font size to 14
    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)

###########
# Save presentation
###########
prs.save('tech_val.pptx')
